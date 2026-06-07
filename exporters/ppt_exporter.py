"""
exporters/ppt_exporter.py
python-pptx 기반 PPT(.pptx) 내보내기 — Node.js 불필요
"""
from __future__ import annotations
import io
from datetime import date

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


# ── 색상 ──────────────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x0D, 0x1B, 0x4B)
BLUE   = RGBColor(0x15, 0x65, 0xC0)
LBLUE  = RGBColor(0x42, 0xA5, 0xF5)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GOLD   = RGBColor(0xFF, 0xD5, 0x4F)
GRAY   = RGBColor(0x54, 0x6E, 0x7A)
LGRAY  = RGBColor(0x90, 0xA4, 0xAE)
RED    = RGBColor(0xEF, 0x53, 0x50)
GREEN  = RGBColor(0x66, 0xBB, 0x6A)
BG     = RGBColor(0xF8, 0xFA, 0xFF)
DARK   = RGBColor(0x1A, 0x23, 0x7E)

W = Inches(10)   # 슬라이드 너비
H = Inches(5.625)  # 슬라이드 높이


# ── 헬퍼 ─────────────────────────────────────────────────────────────────────

def _bg(slide, color: RGBColor):
    """슬라이드 배경색 설정."""
    from pptx.oxml.ns import qn
    from lxml import etree
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _box(slide, x, y, w, h, text="", font_size=18, bold=False,
         color=WHITE, bg_color=None, align=PP_ALIGN.LEFT, italic=False):
    """텍스트 박스 추가."""
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text)
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    if bg_color:
        from pptx.util import Pt as _Pt
        txBox.fill.solid()
        txBox.fill.fore_color.rgb = bg_color
    return txBox


def _rect(slide, x, y, w, h, color: RGBColor):
    """색상 사각형 추가."""
    from pptx.util import Inches as I
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        I(x), I(y), I(w), I(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def _multiline_box(slide, x, y, w, h, lines: list[tuple],
                   default_size=12, default_color=WHITE):
    """
    여러 줄 텍스트박스.
    lines: [(text, bold, font_size, color), ...]
    """
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for (text, bold, fsize, fcolor) in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(2)
        run = p.add_run()
        run.text = str(text)[:200]
        run.font.size = Pt(fsize or default_size)
        run.font.bold = bool(bold)
        run.font.color.rgb = fcolor or default_color
    return txBox


def _change_color(change) -> RGBColor:
    if change is None: return LGRAY
    return RED if change > 0 else LBLUE if change < 0 else LGRAY


def _clean(text: str, limit: int = 200) -> str:
    return text.replace("**", "").replace("##", "").strip()[:limit]


# ── 슬라이드 생성 함수 ────────────────────────────────────────────────────────

def _slide_cover(prs, today, persona_label, persona_emoji):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _bg(slide, NAVY)
    _rect(slide, 0, 3.7, 10, 0.85, BLUE)
    _box(slide, 0.5, 0.5, 9, 1.1,
         "📊 주식 AI 투자 브리핑", font_size=38, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    _box(slide, 0.5, 1.65, 9, 0.5,
         "AI Investment Briefing Report", font_size=17, italic=True, color=LBLUE)
    _box(slide, 0.5, 2.3, 9, 0.55,
         f"{persona_emoji} {persona_label}", font_size=19, bold=True, color=GOLD)
    _box(slide, 0.5, 3.75, 9, 0.5,
         f"기준일: {today}", font_size=15, color=WHITE, align=PP_ALIGN.CENTER)
    _box(slide, 0.5, 5.1, 9, 0.35,
         "본 리포트는 AI 분석 기반이며 투자 결정은 본인 책임입니다.",
         font_size=10, italic=True, color=LGRAY, align=PP_ALIGN.CENTER)


def _slide_portfolio_table(prs, stocks, today):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, BG)
    _rect(slide, 0, 0, 10, 0.85, NAVY)
    _box(slide, 0.3, 0.1, 8, 0.65,
         "📈 포트폴리오 종목 현황", font_size=22, bold=True, color=WHITE)
    _box(slide, 8.0, 0.1, 1.8, 0.65,
         today, font_size=11, color=LGRAY, align=PP_ALIGN.RIGHT)

    # 표 헤더 배경
    cols_x  = [0.3, 3.3, 4.9, 6.5, 8.2]
    cols_w  = [3.0, 1.6, 1.6, 1.7, 1.5]
    headers = ["종목명", "코드/티커", "종가", "등락률", "시장"]
    row_h   = 0.38
    header_y = 0.95

    _rect(slide, 0.3, header_y, 9.5, row_h, BLUE)
    for i, (hdr, cx, cw) in enumerate(zip(headers, cols_x, cols_w)):
        _box(slide, cx, header_y + 0.04, cw, row_h - 0.06,
             hdr, font_size=12, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER if i > 0 else PP_ALIGN.LEFT)

    for row_idx, s in enumerate(stocks[:10]):
        ry = header_y + row_h * (row_idx + 1)
        bg_c = RGBColor(0xE8, 0xF0, 0xFE) if row_idx % 2 == 0 else WHITE
        _rect(slide, 0.3, ry, 9.5, row_h, bg_c)

        name   = s.get("name", s["ticker"])[:18]
        ticker = s["ticker"]
        close  = s.get("close")
        change = s.get("change_rate")
        market = s.get("market", "")
        currency   = "원" if market == "KR" else "$"
        close_str  = f"{close:,.0f}{currency}" if close else "N/A"
        change_str = f"{change:+.2f}%" if change is not None else "N/A"
        clr = _change_color(change)

        vals   = [name, ticker, close_str, change_str, market]
        aligns = [PP_ALIGN.LEFT, PP_ALIGN.CENTER, PP_ALIGN.RIGHT,
                  PP_ALIGN.RIGHT, PP_ALIGN.CENTER]
        colors = [DARK, DARK, DARK, clr, DARK]
        bolds  = [False, False, False, True, False]

        for i, (val, align, color, bold) in enumerate(zip(vals, aligns, colors, bolds)):
            _box(slide, cols_x[i], ry + 0.04, cols_w[i], row_h - 0.06,
                 val, font_size=11, bold=bold, color=color, align=align)


def _slide_strategy(prs, portfolio_brief, persona_label, persona_emoji):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, BG)
    _rect(slide, 0, 0, 10, 0.85, BLUE)
    _box(slide, 0.3, 0.1, 7.5, 0.65,
         "📋 포트폴리오 종합 전략", font_size=22, bold=True, color=WHITE)
    _box(slide, 7.8, 0.1, 2.0, 0.65,
         f"{persona_emoji}", font_size=20, color=GOLD, align=PP_ALIGN.RIGHT)

    # 본문 파싱
    raw_lines = [l.strip() for l in portfolio_brief.split("\n") if l.strip()][:16]
    text_lines = []
    for line in raw_lines:
        is_bold = (line.startswith("**") or line.startswith("##") or
                   (len(line) > 2 and line[1] == "." and line[0].isdigit()))
        clean = _clean(line, 130)
        if clean:
            text_lines.append((clean, is_bold, 13 if is_bold else 12,
                                BLUE if is_bold else DARK))

    _multiline_box(slide, 0.4, 1.0, 9.2, 4.4, text_lines[:14])


def _slide_stock(prs, stock, brief, idx, total, today):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, BG)

    name   = stock.get("name", stock["ticker"])
    ticker = stock["ticker"]
    close  = stock.get("close")
    change = stock.get("change_rate")
    market = stock.get("market", "")
    currency   = "원" if market == "KR" else "$"
    close_str  = f"{close:,.0f}{currency}" if close else "N/A"
    change_str = f"{change:+.2f}%" if change is not None else "N/A"
    flag  = "🇰🇷" if market == "KR" else "🇺🇸"
    clr   = _change_color(change)

    # 헤더 바
    _rect(slide, 0, 0, 10, 0.85, NAVY)
    _box(slide, 0.3, 0.08, 6.5, 0.45,
         f"{flag} {name} ({ticker})", font_size=20, bold=True, color=WHITE)
    _box(slide, 0.3, 0.52, 4, 0.3,
         f"{market}  |  {today}", font_size=11, color=LBLUE)

    # 종가 박스 (우상단)
    _rect(slide, 7.2, 0.08, 2.6, 0.72, BLUE)
    _box(slide, 7.2, 0.08, 2.6, 0.38,
         close_str, font_size=17, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _box(slide, 7.2, 0.46, 2.6, 0.34,
         change_str, font_size=15, bold=True, color=clr, align=PP_ALIGN.CENTER)

    # 좌측: AI 분석 요약
    raw_lines = [l.strip() for l in brief.split("\n") if l.strip()][:12]
    text_lines = []
    for line in raw_lines:
        is_bold = (line.startswith("**") or line.startswith("##") or
                   (len(line) > 2 and line[1] == "." and line[0].isdigit()))
        clean = _clean(line, 110)
        if clean:
            text_lines.append((clean, is_bold, 12 if is_bold else 11,
                                BLUE if is_bold else DARK))

    _box(slide, 0.3, 0.95, 5.2, 0.3,
         "AI 분석 요약", font_size=12, bold=True, color=BLUE)
    _multiline_box(slide, 0.3, 1.28, 5.2, 4.0, text_lines[:11])

    # 우측: 주가 추이 (텍스트 형태)
    hist = stock.get("history", {})
    chart_lines = [("📊 주가 추이", True, 12, BLUE)]
    for period, label in [("week", "이번주"), ("month", "이번달"), ("6month", "6개월")]:
        df = hist.get(period)
        if df is not None and not df.empty and "close" in df.columns:
            sp = float(df["close"].iloc[0])
            ep = float(df["close"].iloc[-1])
            ch = (ep - sp) / sp * 100 if sp else 0
            arrow = "▲" if ch > 0 else "▼" if ch < 0 else "─"
            c = RED if ch > 0 else LBLUE if ch < 0 else LGRAY
            chart_lines.append((f"{label}: {sp:,.0f}→{ep:,.0f} {arrow}{ch:+.1f}%",
                                 False, 11, c))
        else:
            chart_lines.append((f"{label}: 데이터 없음", False, 10, LGRAY))

    _multiline_box(slide, 5.7, 0.95, 4.1, 1.8, chart_lines)

    # 수급 정보
    investor = stock.get("investor", {})
    if investor:
        inv_lines = [("💹 수급 현황", True, 12, BLUE)]
        for k, label in [("외국인합계","외국인"), ("기관합계","기관"), ("개인","개인")]:
            val = investor.get(k)
            if val is not None:
                c = RED if val > 0 else LBLUE
                inv_lines.append((f"{label}: {val:+,}원", False, 10, c))
        _multiline_box(slide, 5.7, 2.9, 4.1, 1.5, inv_lines)

    # 슬라이드 번호
    _box(slide, 0.3, 5.3, 9.4, 0.25,
         f"({idx+1}/{total}) {name}", font_size=9, color=LGRAY, align=PP_ALIGN.RIGHT)


def _slide_ending(prs, persona_label, persona_emoji, today):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, NAVY)
    _box(slide, 0.5, 1.6, 9, 1.2,
         "감사합니다", font_size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _box(slide, 0.5, 2.9, 9, 0.65,
         f"{persona_emoji} {persona_label}", font_size=20, bold=True,
         color=GOLD, align=PP_ALIGN.CENTER)
    _box(slide, 0.5, 4.7, 9, 0.4,
         "본 리포트는 AI 분석 기반이며 투자 결정은 본인 책임입니다.",
         font_size=12, italic=True, color=LGRAY, align=PP_ALIGN.CENTER)
    _box(slide, 0.5, 5.15, 9, 0.3,
         f"기준일: {today}", font_size=11, color=GRAY, align=PP_ALIGN.CENTER)


# ── 메인 내보내기 ────────────────────────────────────────────────────────────

def export_ppt(
    stocks: list[dict],
    stock_briefs: list[str],
    portfolio_brief: str,
    persona_label: str,
    persona_emoji: str,
    output_path: str,
) -> str:
    today = date.today().isoformat()
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(5.625)

    _slide_cover(prs, today, persona_label, persona_emoji)
    _slide_portfolio_table(prs, stocks, today)
    _slide_strategy(prs, portfolio_brief, persona_label, persona_emoji)

    for idx, (stock, brief) in enumerate(zip(stocks, stock_briefs)):
        _slide_stock(prs, stock, brief, idx, len(stocks), today)

    _slide_ending(prs, persona_label, persona_emoji, today)

    prs.save(output_path)
    return output_path


def export_ppt_bytes(
    stocks: list[dict],
    stock_briefs: list[str],
    portfolio_brief: str,
    persona_label: str,
    persona_emoji: str,
) -> bytes:
    """바이트로 반환 (Streamlit download_button용)."""
    import tempfile, os
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
        tmp = f.name
    try:
        export_ppt(stocks, stock_briefs, portfolio_brief,
                   persona_label, persona_emoji, tmp)
        return open(tmp, "rb").read()
    finally:
        os.unlink(tmp)
