"""
uploaders/file_reader.py
업로드 파일 텍스트 추출: Excel, PDF, Word, PPT, TXT, HWP, 이미지
"""
from __future__ import annotations
import os, re, io
from pathlib import Path
from typing import Optional


def extract_text(file_bytes: bytes, filename: str) -> str:
    """파일 바이트 + 파일명으로 텍스트 추출. 실패 시 오류 메시지 반환."""
    ext = Path(filename).suffix.lower()
    try:
        if ext in (".txt", ".md", ".csv"):
            return _read_text(file_bytes)
        elif ext in (".xlsx", ".xls", ".xlsm"):
            return _read_excel(file_bytes)
        elif ext == ".pdf":
            return _read_pdf(file_bytes)
        elif ext in (".docx", ".doc"):
            return _read_word(file_bytes)
        elif ext in (".pptx", ".ppt"):
            return _read_ppt(file_bytes)
        elif ext == ".hwp":
            return _read_hwp(file_bytes)
        elif ext in (".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp"):
            return _read_image(file_bytes, filename)
        else:
            return f"[지원하지 않는 파일 형식: {ext}]"
    except Exception as e:
        return f"[파일 읽기 오류 ({filename}): {type(e).__name__}: {e}]"


# ── TXT ──────────────────────────────────────────────────────────────────────
def _read_text(data: bytes) -> str:
    for enc in ("utf-8", "euc-kr", "cp949", "latin-1"):
        try:
            return data.decode(enc)
        except Exception:
            continue
    return data.decode("utf-8", errors="replace")


# ── Excel ─────────────────────────────────────────────────────────────────────
def _read_excel(data: bytes) -> str:
    import openpyxl
    wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    parts = []
    for shname in wb.sheetnames:
        ws = wb[shname]
        parts.append(f"[시트: {shname}]")
        rows = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(c.strip() for c in cells):
                rows.append(" | ".join(cells))
        parts.extend(rows[:500])  # 최대 500행
    return "\n".join(parts)


# ── PDF ───────────────────────────────────────────────────────────────────────
def _read_pdf(data: bytes) -> str:
    try:
        import pdfplumber
        with pdfplumber.open(io.BytesIO(data)) as pdf:
            pages = []
            for i, page in enumerate(pdf.pages[:30]):  # 최대 30페이지
                txt = page.extract_text() or ""
                if txt.strip():
                    pages.append(f"[p.{i+1}]\n{txt}")
            return "\n\n".join(pages)
    except ImportError:
        # pdfplumber 없으면 PyPDF2 시도
        try:
            import PyPDF2
            reader = PyPDF2.PdfReader(io.BytesIO(data))
            return "\n".join(
                reader.pages[i].extract_text() or ""
                for i in range(min(30, len(reader.pages)))
            )
        except Exception:
            return "[PDF 읽기 실패: pdfplumber 또는 PyPDF2 필요]"


# ── Word (docx) ───────────────────────────────────────────────────────────────
def _read_word(data: bytes) -> str:
    from docx import Document
    doc = Document(io.BytesIO(data))
    paras = [p.text for p in doc.paragraphs if p.text.strip()]
    # 표 내용도 추출
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                paras.append(" | ".join(cells))
    return "\n".join(paras)


# ── PPT (pptx) ────────────────────────────────────────────────────────────────
def _read_ppt(data: bytes) -> str:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(data))
    slides_text = []
    for i, slide in enumerate(prs.slides):
        parts = [f"[슬라이드 {i+1}]"]
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text.strip())
        slides_text.append("\n".join(parts))
    return "\n\n".join(slides_text)


# ── HWP (한글) ────────────────────────────────────────────────────────────────
def _read_hwp(data: bytes) -> str:
    """HWP는 olefile 기반 간이 추출 (완전하지 않을 수 있음)."""
    try:
        import olefile
        f = olefile.OleFileIO(io.BytesIO(data))
        if f.exists("PrvText"):
            raw = f.openstream("PrvText").read()
            return raw.decode("utf-16-le", errors="replace")
        # BodyText 스트림에서 직접 추출
        texts = []
        for entry in f.listdir():
            name = "/".join(entry)
            if "BodyText" in name:
                try:
                    raw = f.openstream(entry).read()
                    txt = re.sub(rb"[\x00-\x08\x0b\x0c\x0e-\x1f]", b"", raw)
                    decoded = txt.decode("utf-16-le", errors="replace")
                    if decoded.strip():
                        texts.append(decoded)
                except Exception:
                    pass
        return "\n".join(texts) if texts else "[HWP: 텍스트 추출 실패]"
    except ImportError:
        return "[HWP 읽기 실패: olefile 패키지 필요 → pip install olefile]"


# ── 이미지 (OCR 또는 base64 → LLM) ──────────────────────────────────────────
def _read_image(data: bytes, filename: str) -> str:
    """이미지에서 텍스트 추출. pytesseract 있으면 OCR, 없으면 안내."""
    try:
        import pytesseract
        from PIL import Image
        img = Image.open(io.BytesIO(data))
        text = pytesseract.image_to_string(img, lang="kor+eng")
        return f"[이미지 OCR: {filename}]\n{text}" if text.strip() else f"[이미지 텍스트 없음: {filename}]"
    except ImportError:
        # OCR 없을 때 base64로 저장해 LLM 분석용 메타데이터 반환
        import base64
        b64 = base64.b64encode(data).decode()
        ext = Path(filename).suffix.lower().replace(".", "")
        media_type = {"jpg": "image/jpeg", "jpeg": "image/jpeg",
                      "png": "image/png", "gif": "image/gif"}.get(ext, "image/png")
        return f"[IMAGE_B64:{media_type}:{b64[:100]}...]"  # 미리보기만


def summarize_file_content(text: str, filename: str, max_chars: int = 3000) -> str:
    """긴 파일 내용을 RAG 청크용으로 트리밍."""
    if len(text) <= max_chars:
        return text
    return text[:max_chars] + f"\n...[{filename}: 총 {len(text)}자 중 {max_chars}자 표시]"
